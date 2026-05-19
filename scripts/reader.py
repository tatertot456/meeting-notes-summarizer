import fitz  # PyMuPDF
from docx import Document
from pathlib import Path
from pptx import Presentation


def read_file(filepath: str) -> str:
    path = Path(filepath)
    suffix = path.suffix.lower()

    if suffix == ".txt" or suffix == ".md":
        return path.read_text(encoding="utf-8")

    elif suffix == ".pdf":
        doc = fitz.open(filepath)
        return "\n".join(page.get_text() for page in doc)

    elif suffix == ".docx":
        doc = Document(filepath)
        return "\n".join(p.text for p in doc.paragraphs)

    elif suffix == ".pptx":
        prs = Presentation(filepath)
        slides_text = []
        for i, slide in enumerate(prs.slides):
            slide_lines = [f"--- Slide {i + 1} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_lines.append(text)
            slides_text.append("\n".join(slide_lines))
        return "\n\n".join(slides_text)

    else:
        raise ValueError(f"Unsupported file type: {suffix}")